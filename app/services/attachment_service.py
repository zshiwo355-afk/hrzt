"""附件上传、文本抽取、对话上下文中的附件描述。"""
from __future__ import annotations

import base64
import io
import tempfile
import re
import subprocess
import shutil
import uuid
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any, Optional

from fastapi import HTTPException, UploadFile

from app.config import ALLOWED_EXTENSIONS, OSS_UPLOAD_PREFIX, UPLOADS_DIR, UPLOAD_META_DIR, oss_configured
from app.database_models import Attachment
from app.db import get_session_factory
from app.logging_config import logger
from app.providers import oss as oss_provider
from app.storage import read_attachment_meta, safe_json_dump

THUMBNAILS_DIR = UPLOADS_DIR / "_thumbs"
THUMBNAIL_MAX_LONG_EDGE = 768
THUMBNAIL_QUALITY = 82


def truncate_text(text: str, limit: int = 12000) -> str:
    text = (text or "").strip()
    if len(text) <= limit:
        return text
    return text[:limit] + "\n\n[内容过长，已截断]"


def detect_file_category(filename: str, mime_type: str) -> str:
    suffix = Path(filename).suffix.lower()

    if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
        return "image"

    if suffix in {".pdf", ".doc", ".docx", ".ppt", ".pptx", ".txt", ".md", ".xlsx", ".xls"}:
        return "document"

    if mime_type.startswith("image/"):
        return "image"

    return "other"


async def save_upload_file(upload: UploadFile) -> dict:
    original_name = Path(upload.filename or "未命名文件").name
    suffix = Path(original_name).suffix.lower()

    if suffix not in ALLOWED_EXTENSIONS:
        raise HTTPException(status_code=400, detail=f"暂不支持上传这种文件类型：{original_name}")

    attachment_id = uuid.uuid4().hex
    stored_name = f"{attachment_id}{suffix}"
    try:
        content = await upload.read()
    finally:
        await upload.close()

    mime_type = upload.content_type or "application/octet-stream"
    category = detect_file_category(original_name, mime_type)
    size = len(content or b"")
    oss_key = ""
    stored_path = UPLOADS_DIR / stored_name
    storage = "local"
    if oss_configured():
        oss_key = f"{OSS_UPLOAD_PREFIX}/{attachment_id}{suffix}"
        if oss_provider.write_bytes(oss_key, content, content_type=mime_type):
            storage = "oss"
    with stored_path.open("wb") as f:
        f.write(content)

    meta = {
        "id": attachment_id,
        "original_name": original_name,
        "stored_name": stored_name,
        "stored_path": str(stored_path),
        "storage": storage,
        "oss_key": oss_key if storage == "oss" else "",
        "suffix": suffix,
        "mime_type": mime_type,
        "category": category,
        "size": size,
    }

    safe_json_dump(UPLOAD_META_DIR / f"{attachment_id}.json", meta)
    save_attachment_meta_db(meta, source="upload")
    return meta


def save_attachment_meta_db(
    meta: dict,
    *,
    user_id: Optional[int] = None,
    conversation_id: Optional[int] = None,
    message_id: Optional[int] = None,
    task_id: Optional[str] = None,
    source: str = "upload",
) -> None:
    if not meta or not meta.get("id"):
        return
    try:
        with get_session_factory()() as db:
            row = db.get(Attachment, str(meta["id"]))
            if row is None:
                row = Attachment(id=str(meta["id"]))
            row.user_id = int(user_id) if user_id else row.user_id
            row.conversation_id = int(conversation_id) if conversation_id else row.conversation_id
            row.message_id = int(message_id) if message_id else row.message_id
            row.task_id = str(task_id) if task_id else row.task_id
            row.original_name = str(meta.get("original_name") or "")
            row.stored_name = str(meta.get("stored_name") or "")
            row.stored_path = str(meta.get("stored_path") or "")
            row.storage = str(meta.get("storage") or "local")
            row.oss_key = str(meta.get("oss_key") or "")
            row.suffix = str(meta.get("suffix") or "")
            row.mime_type = str(meta.get("mime_type") or "application/octet-stream")
            row.category = str(meta.get("category") or "document")
            row.size = int(meta.get("size") or 0)
            row.source = str(source or "upload")
            db.add(row)
            db.commit()
    except Exception as exc:
        logger.warning("[attachment-db] save meta failed attachment_id=%s err=%r", meta.get("id"), exc)


def read_attachment_meta_any(attachment_id: str) -> Optional[dict]:
    meta = read_attachment_meta(attachment_id)
    if meta:
        return meta
    try:
        with get_session_factory()() as db:
            row = db.get(Attachment, str(attachment_id))
            if not row:
                return None
            return {
                "id": row.id,
                "original_name": row.original_name,
                "stored_name": row.stored_name,
                "stored_path": row.stored_path,
                "storage": row.storage,
                "oss_key": row.oss_key,
                "suffix": row.suffix,
                "mime_type": row.mime_type,
                "category": row.category,
                "size": row.size,
            }
    except Exception as exc:
        logger.warning("[attachment-db] read meta failed attachment_id=%s err=%r", attachment_id, exc)
        return None


def build_attachment_public(meta: dict) -> dict:
    return {
        "id": meta["id"],
        "name": meta["original_name"],
        "mime_type": meta["mime_type"],
        "category": meta["category"],
        "size": meta["size"],
    }


def read_attachment_bytes(meta: dict) -> Optional[bytes]:
    oss_key = str(meta.get("oss_key") or "").strip()
    if oss_key:
        data = oss_provider.read_bytes(oss_key)
        if data is not None:
            return data
    stored_path = str(meta.get("stored_path") or "").strip()
    if not stored_path:
        return None
    path = Path(stored_path)
    if not path.is_file():
        return None
    return path.read_bytes()


def build_image_thumbnail_bytes(meta: dict) -> tuple[bytes, str] | None:
    if not meta or str(meta.get("category") or "") != "image":
        return None

    attachment_id = str(meta.get("id") or "").strip()
    if not attachment_id:
        return None

    THUMBNAILS_DIR.mkdir(parents=True, exist_ok=True)
    thumb_path = THUMBNAILS_DIR / f"{attachment_id}.webp"
    if thumb_path.is_file():
        try:
            return thumb_path.read_bytes(), "image/webp"
        except Exception:
            pass

    raw = None
    stored_path = str(meta.get("stored_path") or "").strip()
    if stored_path:
        path = Path(stored_path)
        if path.is_file():
            try:
                raw = path.read_bytes()
            except OSError:
                raw = None
    if raw is None:
        raw = read_attachment_bytes(meta)
    if raw is None:
        return None

    try:
        from PIL import Image, ImageOps

        im = Image.open(io.BytesIO(raw))
        im = ImageOps.exif_transpose(im)
        im.thumbnail((THUMBNAIL_MAX_LONG_EDGE, THUMBNAIL_MAX_LONG_EDGE), Image.Resampling.LANCZOS)
        if im.mode not in ("RGB", "RGBA"):
            im = im.convert("RGB")

        out = io.BytesIO()
        save_kwargs = {"format": "WEBP", "quality": THUMBNAIL_QUALITY, "method": 4}
        if im.mode == "RGBA":
            save_kwargs["lossless"] = False
        im.save(out, **save_kwargs)
        data = out.getvalue()
        try:
            thumb_path.write_bytes(data)
        except Exception:
            logger.warning("[thumb] cache write failed attachment_id=%s", attachment_id, exc_info=True)
        return data, "image/webp"
    except Exception:
        logger.warning("[thumb] build failed attachment_id=%s", attachment_id, exc_info=True)
        return None


def materialize_attachment_file(meta: dict) -> Optional[Path]:
    stored_path = str(meta.get("stored_path") or "").strip()
    if stored_path:
        path = Path(stored_path)
        if path.is_file():
            return path
    data = read_attachment_bytes(meta)
    if data is None:
        return None
    suffix = str(meta.get("suffix") or "").strip() or ".bin"
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    try:
        tmp.write(data)
        tmp.flush()
        return Path(tmp.name)
    finally:
        tmp.close()


# Backward-compatible alias for any long-lived worker that imports the older helper name.
materialize_attachment = materialize_attachment_file


def extract_docx_text(file_path: Path) -> str:
    try:
        from docx import Document

        doc = Document(str(file_path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        table_rows = []
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    table_rows.append(" | ".join(cells))

        parts = paragraphs + table_rows
        return "\n\n".join(parts).strip()
    except ImportError:
        pass

    try:
        with zipfile.ZipFile(file_path) as zf:
            xml_bytes = zf.read("word/document.xml")
        root = ET.fromstring(xml_bytes)

        ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
        paragraphs = []
        for p_node in root.iter(f"{ns}p"):
            runs = []
            for t_node in p_node.iter(f"{ns}t"):
                if t_node.text:
                    runs.append(t_node.text)
            line = "".join(runs).strip()
            if line:
                paragraphs.append(line)

        return "\n\n".join(paragraphs).strip()
    except Exception:
        return ""


def extract_doc_text(file_path: Path) -> str:
    if shutil.which("antiword"):
        try:
            result = subprocess.run(
                ["antiword", str(file_path)],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip()
        except Exception:
            pass

    if shutil.which("catdoc"):
        try:
            result = subprocess.run(
                ["catdoc", str(file_path)],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip()
        except Exception:
            pass

    try:
        raw = file_path.read_bytes()
        text_chunks = []
        i = 0
        while i < len(raw):
            start = i
            while i < len(raw) and 0x20 <= raw[i] < 0x7F:
                i += 1
            if i - start >= 10:
                text_chunks.append(raw[start:i].decode("ascii", errors="ignore"))
            i += 1
        if text_chunks:
            return "\n".join(text_chunks).strip()
    except Exception:
        pass

    return "[.doc 文件已上传。服务器未安装 antiword/catdoc，无法自动抽取旧版 Word 正文。建议转存为 .docx 后重新上传。]"


def extract_xlsx_text(file_path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except Exception:
        return "[Excel 已上传，但服务器暂未安装 openpyxl，无法抽取内容。]"

    try:
        wb = load_workbook(str(file_path), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append("\t".join(cells))
            if rows:
                parts.append(f"【工作表：{sheet_name}】\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts).strip() if parts else "[Excel 已上传，但未读取到有效内容。]"
    except Exception:
        return "[Excel 已上传，但读取内容失败。]"


def extract_pdf_text(file_path: Path) -> str:
    try:
        from pypdf import PdfReader
    except Exception:
        return "[PDF 已上传，但服务器暂未启用 PDF 正文抽取。]"

    try:
        reader = PdfReader(str(file_path))
        parts = []
        for page in reader.pages:
            text = page.extract_text() or ""
            text = text.strip()
            if text:
                parts.append(text)
        return "\n\n".join(parts).strip()
    except Exception:
        return "[PDF 已上传，但读取正文失败。]"


def extract_pptx_text(file_path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception:
        return "[PPTX 已上传，但服务器暂未安装 python-pptx，无法抽取内容。]"

    try:
        prs = Presentation(str(file_path))
        slides = []
        for idx, slide in enumerate(prs.slides, start=1):
            parts = []
            for shape in slide.shapes:
                text = ""
                if getattr(shape, "has_text_frame", False) and shape.text_frame:
                    text = "\n".join(
                        p.text.strip()
                        for p in shape.text_frame.paragraphs
                        if p.text and p.text.strip()
                    )
                elif getattr(shape, "has_table", False) and shape.table:
                    rows = []
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            rows.append(" | ".join(cells))
                    text = "\n".join(rows)
                if text:
                    parts.append(text)
            if getattr(slide, "has_notes_slide", False) and slide.notes_slide:
                notes = []
                for shape in slide.notes_slide.shapes:
                    if getattr(shape, "has_text_frame", False) and shape.text_frame:
                        for p in shape.text_frame.paragraphs:
                            t = (p.text or "").strip()
                            if t and t.lower() != "slide notes":
                                notes.append(t)
                if notes:
                    parts.append("备注：\n" + "\n".join(notes))
            if parts:
                slides.append(f"【第 {idx} 页】\n" + "\n".join(parts))
        return "\n\n".join(slides).strip() if slides else "[PPTX 已上传，但未读取到有效文本。]"
    except Exception:
        return "[PPTX 已上传，但读取内容失败。]"


def extract_ppt_text(file_path: Path) -> str:
    return "[.ppt 文件已上传。旧版 PPT 为二进制格式，当前服务器暂未启用自动解析；建议另存为 .pptx 后可读取幻灯片正文。]"


def extract_text_from_attachment_meta(meta: dict) -> str:
    suffix = (meta.get("suffix") or "").lower()
    file_path = materialize_attachment_file(meta)
    if not file_path:
        return ""

    if suffix in {".txt", ".md"}:
        try:
            return file_path.read_text(encoding="utf-8", errors="ignore").strip()
        except Exception:
            return ""

    if suffix == ".docx":
        return extract_docx_text(file_path)

    if suffix == ".pdf":
        return extract_pdf_text(file_path)

    if suffix == ".doc":
        return extract_doc_text(file_path)

    if suffix == ".pptx":
        return extract_pptx_text(file_path)

    if suffix == ".ppt":
        return extract_ppt_text(file_path)

    if suffix in {".xlsx", ".xls"}:
        return extract_xlsx_text(file_path)

    return ""


def build_current_attachment_context(attachment_ids: list[str]) -> str:
    if not attachment_ids:
        return ""

    docs = []
    images = []
    notes = []

    for attachment_id in attachment_ids:
        meta = read_attachment_meta_any(attachment_id)
        if not meta:
            continue

        name = meta["original_name"]
        category = meta["category"]

        if category == "image":
            images.append(name)
            continue

        extracted = extract_text_from_attachment_meta(meta)
        extracted = truncate_text(extracted, 12000)

        if extracted:
            docs.append(f"文件名：{name}\n文件内容如下：\n{extracted}")
        else:
            notes.append(f"{name} 已上传，但暂未解析出可用文本内容。")

    sections = []

    if docs:
        sections.append("以下是当前上传的文件内容：\n\n" + "\n\n===== 分隔 =====\n\n".join(docs))

    if images:
        sections.append("当前还上传了这些图片文件：" + "、".join(images))

    if notes:
        sections.append("附件说明：" + "；".join(notes))

    return "\n\n".join(sections).strip()


def build_attachment_names(attachment_ids: list[str], category: Optional[str] = None) -> list[str]:
    names = []
    for attachment_id in attachment_ids:
        meta = read_attachment_meta_any(attachment_id)
        if not meta:
            continue
        if category and meta.get("category") != category:
            continue
        names.append(meta["original_name"])
    return names


def extract_text_from_content(content: Any) -> str:
    if content is None:
        return ""

    if isinstance(content, str):
        return content

    if isinstance(content, list):
        parts = []
        for item in content:
            if isinstance(item, str):
                parts.append(item)
                continue

            if isinstance(item, dict):
                item_type = str(item.get("type", "")).lower()

                if item_type in {"text", "output_text", "message"}:
                    text_value = item.get("text") or item.get("content") or item.get("value") or ""
                    if isinstance(text_value, str):
                        parts.append(text_value)
                    elif isinstance(text_value, list):
                        parts.append(extract_text_from_content(text_value))
                    continue

                if "text" in item and isinstance(item["text"], str):
                    parts.append(item["text"])
                    continue

                if "content" in item:
                    parts.append(extract_text_from_content(item.get("content")))
                    continue

        return "\n".join(p for p in parts if p).strip()

    if isinstance(content, dict):
        if isinstance(content.get("text"), str):
            return content["text"]

        if "content" in content:
            return extract_text_from_content(content.get("content"))

        if isinstance(content.get("message"), dict):
            return extract_text_from_content(content["message"].get("content"))

    return str(content).strip()


def clean_response_text(text: str) -> str:
    if not text:
        return ""

    cleaned = text.replace("\r\n", "\n").strip()

    noisy_prefixes = [
        r"^好的[，,。:\s]*",
        r"^好[，,。:\s]*",
        r"^当然可以[，,。:\s]*",
        r"^下面[我先]*来[给你]*[一二三四五六七八九十0-9]*[点步]*[分析说明整理概括展开回答]*[：:\s]*",
        r"^我来[先]*[给你]*[一二三四五六七八九十0-9]*[点步]*[分析说明整理展开回答]*[：:\s]*",
        r"^先来[看说讲][一下一点些]*[：:\s]*",
        r"^我们先[来看想说讲分析一下整理一下]*[：:\s]*",
        r"^先[分析说明整理回答][一下一点些]*[：:\s]*",
        r"^我的思路是[：:\s]*",
        r"^我先说结论[：:\s]*",
        r"^先说结论[：:\s]*",
        r"^让我先[想分析看]一下[：:\s]*",
        r"^我先[想分析]一下[：:\s]*",
        r"^以下是[我给你的]*[回答分析建议内容][：:\s]*",
    ]

    for pattern in noisy_prefixes:
        cleaned = re.sub(pattern, "", cleaned, flags=re.IGNORECASE)

    chain_patterns = [
        r"<think>.*?</think>",
        r"<thinking>.*?</thinking>",
        r"```thinking[\s\S]*?```",
        r"```reasoning[\s\S]*?```",
    ]
    for pattern in chain_patterns:
        cleaned = re.sub(pattern, "", cleaned, flags=re.IGNORECASE)

    lines = [line.rstrip() for line in cleaned.split("\n")]
    filtered_lines = []
    skip_head_noise = True
    skip_process_paragraph = False
    skipped_process_blocks = 0
    process_heading_re = re.compile(
        r"^\s*(?:#{1,6}\s*)?"
        r"(?:[-*]\s*)?"
        r"(?:\*\*|__)?"
        r"(?:"
        r"Analy[sz]ing|Exploring|Retrieving|Fetching|Searching|Investigating|Validating|Checking|"
        r"Reviewing|Identifying|Processing|Gathering|Reading|Scanning|Looking\s+up|"
        r"Preparing|Understanding|Planning|Thinking|Reasoning|Finding|Discovering|Recommending|"
        r"Selecting|Curating|Comparing|Ranking|Evaluating|Assessing|Examining|Sifting|Shortlisting"
        r")\b.*$",
        flags=re.IGNORECASE,
    )
    process_sentence_re = re.compile(
        r"^\s*(?:This is your|This week|I'?m|I am|I'?ve|I have|I will|I'll|Let me|"
        r"Now I(?:'m| am)|My focus is|My aim is|This suggests|I need to|We're|We are|These include)\b.*"
        r"(?:search|process|analy[sz]|fetch|dig|investigat|validat|check|look|"
        r"identify|identified|focus|drill|understand|provide|bypass|adapt|organize|gather|retriev|"
        r"checking|trending|repositories|infrastructure|rewrite|movement|leading|projects|innovative|"
        r"discussions|illustrat|include|pushing|embodied)",
        flags=re.IGNORECASE,
    )
    likely_answer_start_re = re.compile(r"^\s*(?:[\u4e00-\u9fff]|#{1,6}\s*[\u4e00-\u9fff]|[-*]\s*[\u4e00-\u9fff]|\d+[.)、]\s*[\u4e00-\u9fff])")

    for line in lines:
        stripped = line.strip()

        if skip_head_noise:
            if skip_process_paragraph:
                if not stripped:
                    skip_process_paragraph = False
                elif likely_answer_start_re.match(stripped):
                    filtered_lines.append(line)
                    skip_head_noise = False
                    skip_process_paragraph = False
                continue

            if not stripped:
                continue

            if process_heading_re.match(stripped):
                skip_process_paragraph = True
                skipped_process_blocks += 1
                continue

            if skipped_process_blocks and process_sentence_re.match(stripped):
                continue
            if not skipped_process_blocks and process_sentence_re.match(stripped):
                skipped_process_blocks += 1
                continue

            head_noise_markers = [
                "好的，下面",
                "好的，先",
                "下面我来",
                "我来一步步",
                "我来一步一步",
                "先分析一下",
                "先说一下",
                "我的分析如下",
                "我的思考如下",
                "思路如下",
                "我们来分析一下",
            ]
            if any(stripped.startswith(marker) for marker in head_noise_markers):
                continue

        filtered_lines.append(line)
        if stripped:
            skip_head_noise = False

    cleaned = "\n".join(filtered_lines).strip()
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)

    if not cleaned:
        cleaned = "模型没有返回可显示内容。"

    return cleaned
