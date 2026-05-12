"""内置 PPT 模板渲染：优先复用现成模板，失败时由上层回退固定布局。"""
from __future__ import annotations

import io
import json
import re
from copy import deepcopy
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Emu, Pt


_ROOT_DIR = Path(__file__).resolve().parents[2]
_TEMPLATE_BASE_DIR = _ROOT_DIR / "data" / "ppt_templates"
_TEMPLATE_LIBRARY_FILE = _TEMPLATE_BASE_DIR / "library.json"
_SUPPORTED_TEMPLATE_TYPES = {"cover", "agenda", "content", "ending"}

_PLACEHOLDER_PATTERNS = [
    r"加入你的",
    r"加入标题",
    r"加入文字",
    r"点击添加",
    r"点击输入",
    r"点击此处",
    r"这里输入",
    r"在这里输入",
    r"添加标题",
    r"添加文字",
    r"添加副标题",
    r"添加内容",
    r"您的文字",
    r"您的内容",
    r"lorem ipsum",
    r"click to add",
]
_PLACEHOLDER_RE = re.compile("|".join(_PLACEHOLDER_PATTERNS), re.IGNORECASE)
_DECORATION_RE = re.compile(
    r"^("
    r"\d{1,4}|"
    r"20[0-9x]{2}|"
    r"[A-Za-z]|"
    r"VS|vs|Vs|"
    r"[•●▶▪◆\\-—_]+|"
    r"PPT|TEMPLATE|"
    r"[一二三四五六七八九十]"
    r")$"
)


def is_template_eligible(schema: dict) -> bool:
    slides = schema.get("slides") or []
    if not slides:
        return False
    return all(str(slide.get("type") or "").strip().lower() in _SUPPORTED_TEMPLATE_TYPES for slide in slides)


def list_template_library() -> list[dict]:
    payload = _load_template_library()
    return list(payload.get("templates") or [])


def render_with_builtin_template(schema: dict, *, prompt: str = "") -> tuple[bytes, dict[str, str]]:
    template_bundle = _select_template_bundle(schema, prompt=prompt)
    layout = json.loads((template_bundle["path"] / "template.json").read_text(encoding="utf-8"))
    source_bytes = (template_bundle["path"] / "source.pptx").read_bytes()
    slides_payload = _build_template_payload(schema, layout)
    content = render_pptx_from_slides(
        source_bytes,
        slides_payload,
        overall_title=str(schema.get("title") or "演示文稿"),
        template_layout=layout,
    )
    return content, {
        "template_id": str(template_bundle["id"]),
        "template_label": str(template_bundle["label"]),
        "template_name": str(template_bundle["name"]),
    }


def _load_template_library() -> dict:
    if _TEMPLATE_LIBRARY_FILE.is_file():
        payload = json.loads(_TEMPLATE_LIBRARY_FILE.read_text(encoding="utf-8"))
        if isinstance(payload, dict) and isinstance(payload.get("templates"), list):
            return payload
    return {"version": 1, "templates": []}


def _build_selection_text(schema: dict, prompt: str = "") -> str:
    return " ".join(
        [
            str(prompt or ""),
            str(schema.get("theme_hint") or ""),
            str(schema.get("title") or ""),
            str(schema.get("subtitle") or ""),
            " ".join(str(slide.get("title") or "") for slide in schema.get("slides") or []),
            " ".join(" ".join(str(item) for item in (slide.get("bullets") or [])) for slide in schema.get("slides") or []),
            " ".join(" ".join(str(item) for item in (slide.get("paragraphs") or [])) for slide in schema.get("slides") or []),
        ]
    ).lower()


def _select_template_bundle(schema: dict, *, prompt: str = "") -> dict[str, Any]:
    text = _build_selection_text(schema, prompt=prompt)
    candidates: list[dict[str, Any]] = []
    for template in list_template_library():
        path = _TEMPLATE_BASE_DIR / str(template.get("id") or "")
        if not ((path / "source.pptx").is_file() and (path / "template.json").is_file()):
            continue
        score = _score_template(template, text)
        candidates.append({**template, "path": path, "score": score})
    if candidates:
        candidates.sort(key=lambda item: (item.get("score", 0), item.get("priority", 0)), reverse=True)
        return candidates[0]
    raise FileNotFoundError("未找到可用的内置 PPT 模板资源")


def _score_template(template: dict, text: str) -> int:
    score = int(template.get("priority") or 0)
    label = str(template.get("label") or "").lower()
    name = str(template.get("name") or "").lower()
    desc = str(template.get("description") or "").lower()
    style = str(template.get("style") or "").lower()
    if label and label in text:
        score += 50
    if name and name in text:
        score += 60
    for token in template.get("tags") or []:
        normalized = str(token).strip().lower()
        if normalized and normalized in text:
            score += 18
    for token in template.get("scenes") or []:
        normalized = str(token).strip().lower()
        if normalized and normalized in text:
            score += 24
    for token in filter(None, [desc, style]):
        if token and any(part and part in text for part in re.split(r"[\s,，/]+", token)):
            score += 8
    if any(token in text for token in ("汇报", "周报", "月报", "经营", "述职", "总结")) and "gold" in label:
        score += 25
    if any(token in text for token in ("培训", "项目", "产品", "方案", "深色", "专业")) and "dark" in label:
        score += 25
    return score


def _build_template_payload(schema: dict, layout: dict) -> list[dict]:
    slides = schema.get("slides") or []
    template_indices = _assign_template_indices(slides, layout)
    payload: list[dict] = []
    for idx, slide in enumerate(slides):
        slide_type = str(slide.get("type") or "content").strip().lower()
        bullets = [str(item).strip() for item in (slide.get("bullets") or []) if str(item).strip()]
        paragraphs = [str(item).strip() for item in (slide.get("paragraphs") or []) if str(item).strip()]
        subtitle = str(slide.get("subtitle") or "").strip()
        body = "\n\n".join(paragraphs)
        if slide_type == "agenda" and not bullets and paragraphs:
            bullets = paragraphs
            body = ""
        payload.append(
            {
                "template_page_idx": template_indices[idx],
                "title": str(slide.get("title") or "").strip(),
                "subtitle": subtitle,
                "bullets": bullets,
                "body": body,
                "speaker_note": str(slide.get("notes") or "").strip(),
            }
        )
    return payload


def _assign_template_indices(slides: list[dict], layout: dict) -> list[int]:
    layout_slides = layout.get("slides") or []
    cover_indices: list[int] = []
    end_indices: list[int] = []
    content_indices: list[int] = []
    section_indices: list[int] = []
    for idx, slide in enumerate(layout_slides):
        kind = str(slide.get("layout_kind") or slide.get("page_type") or "").strip().lower()
        if kind == "cover":
            cover_indices.append(idx)
        elif kind == "end":
            end_indices.append(idx)
        elif kind == "section":
            section_indices.append(idx)
        elif kind == "content":
            content_indices.append(idx)
    if not content_indices:
        content_indices = section_indices[:]
    if not content_indices:
        start = 1 if len(layout_slides) > 1 else 0
        end_limit = max(start + 1, len(layout_slides) - 1)
        content_indices = list(range(start, end_limit))
    if not content_indices:
        content_indices = [0]

    assigned: list[int] = []
    content_cursor = 0
    for slide in slides:
        slide_type = str(slide.get("type") or "content").strip().lower()
        if slide_type == "cover":
            assigned.append(cover_indices[0] if cover_indices else 0)
            continue
        if slide_type == "ending":
            assigned.append(end_indices[0] if end_indices else content_indices[-1])
            continue
        assigned.append(content_indices[content_cursor % len(content_indices)])
        content_cursor += 1
    return assigned


def _emu_to_pt(value: Any) -> float:
    if value is None:
        return 0.0
    try:
        return round(float(Emu(value).pt), 2)
    except Exception:
        return 0.0


def _shape_avg_font_size(shape) -> float:
    if not shape.has_text_frame:
        return 0.0
    sizes: list[float] = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size:
                sizes.append(run.font.size.pt)
    if not sizes:
        return 0.0
    return sum(sizes) / len(sizes)


def _shape_text_length(shape) -> int:
    if not shape.has_text_frame:
        return 0
    return len((shape.text_frame.text or "").strip())


def _is_decoration_text(text: str) -> bool:
    cleaned = (text or "").strip()
    if not cleaned:
        return True
    return len(cleaned) <= 6 and bool(_DECORATION_RE.match(cleaned))


def _is_placeholder_text(text: str) -> bool:
    cleaned = (text or "").strip()
    if not cleaned:
        return False
    return bool(_PLACEHOLDER_RE.search(cleaned))


def _bucket_signature(shape) -> tuple[int, int, int]:
    width_pt = round(_emu_to_pt(shape.width) / 20) * 20
    height_pt = round(_emu_to_pt(shape.height) / 10) * 10
    font = round(_shape_avg_font_size(shape))
    return (width_pt, height_pt, font)


def _pick_title_shape(text_shapes: list, role_hint: object | None) -> object | None:
    if role_hint is not None:
        return role_hint
    for shape in text_shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 0:
            return shape
    candidates = [shape for shape in text_shapes if not _is_decoration_text(shape.text_frame.text or "")]
    if not candidates:
        candidates = list(text_shapes)
    candidates.sort(key=lambda shape: (-_shape_avg_font_size(shape), shape.top if shape.top is not None else 0))
    return candidates[0] if candidates else None


def _pick_body_slots(text_shapes: list, title_shape, role_body_hints: list) -> list:
    if role_body_hints:
        return list(role_body_hints)
    title_size = _shape_avg_font_size(title_shape) if title_shape else 0
    candidates = []
    for shape in text_shapes:
        if shape is title_shape:
            continue
        text = shape.text_frame.text or ""
        if _is_decoration_text(text):
            continue
        if title_size and _shape_avg_font_size(shape) > title_size + 1:
            continue
        candidates.append(shape)
    if not candidates:
        return []
    placeholder_shapes = [shape for shape in candidates if _is_placeholder_text(shape.text_frame.text)]
    pool = placeholder_shapes if placeholder_shapes else candidates
    buckets: dict[tuple[int, int, int], list] = {}
    for shape in pool:
        buckets.setdefault(_bucket_signature(shape), []).append(shape)
    if buckets:
        best = max(buckets.values(), key=len)
        if len(best) >= 2:
            best.sort(key=lambda shape: (shape.top if shape.top is not None else 0, shape.left if shape.left is not None else 0))
            return best
    fallback = max(pool, key=_shape_text_length) if pool else None
    return [fallback] if fallback else []


def _classify_slide_shapes(slide, role_by_index: dict[int, str] | None = None) -> tuple[object | None, list, list]:
    text_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
    if not text_shapes:
        return None, [], []
    title_hint = None
    body_hints: list = []
    if role_by_index:
        all_shapes = list(slide.shapes)
        for idx, shape in enumerate(all_shapes):
            role = role_by_index.get(idx)
            if not role or not shape.has_text_frame:
                continue
            if role == "title" and title_hint is None:
                title_hint = shape
            elif role in {"subtitle", "bullets", "content", "item"}:
                body_hints.append(shape)
    title_shape = _pick_title_shape(text_shapes, title_hint)
    body_shapes = _pick_body_slots(text_shapes, title_shape, body_hints)
    used = {id(title_shape)} | {id(shape) for shape in body_shapes}
    leftover = [shape for shape in text_shapes if id(shape) not in used]
    return title_shape, body_shapes, leftover


def _enable_text_autofit(text_frame) -> None:
    try:
        text_frame.word_wrap = True
        from lxml import etree

        body_pr = text_frame._txBody.find("{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr")
        if body_pr is None:
            return
        nsmap = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        for tag in ("spAutoFit", "noAutofit", "normAutofit"):
            existing = body_pr.find(f"{nsmap}{tag}")
            if existing is not None:
                body_pr.remove(existing)
        etree.SubElement(body_pr, f"{nsmap}normAutofit")
        body_pr.set("wrap", "square")
    except Exception:
        pass


def _shrink_font_if_overflow(shape, new_text: str, *, min_pt: float = 8.0) -> None:
    if not new_text or not shape.has_text_frame:
        return
    try:
        width_pt = _emu_to_pt(shape.width)
        height_pt = _emu_to_pt(shape.height)
        if width_pt <= 0 or height_pt <= 0:
            return
        text_len = len(new_text)
        if text_len <= 0:
            return
        current_size = 0.0
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size:
                    current_size = run.font.size.pt
                    break
            if current_size:
                break
        if current_size <= 0:
            return
        capacity = (width_pt / current_size) * max(1.0, height_pt / (current_size * 1.3))
        if capacity >= text_len * 0.95:
            return
        import math

        target_size = math.sqrt(width_pt * height_pt / (text_len * 0.95 * 1.3))
        target_size = min(current_size, max(min_pt, target_size))
        if target_size >= current_size - 0.1:
            return
        new_size = Pt(round(target_size, 1))
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size:
                    run.font.size = new_size
    except Exception:
        pass


def _replace_text_keeping_format(text_frame, new_text: str) -> None:
    if not new_text:
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = ""
        return
    paragraphs = text_frame.paragraphs
    if not paragraphs:
        text_frame.text = new_text
        return
    first_para = paragraphs[0]
    runs = list(first_para.runs)
    if runs:
        runs[0].text = new_text
        for run in runs[1:]:
            run.text = ""
    else:
        first_para.add_run().text = new_text
    for el in [para._p for para in paragraphs[1:]]:
        el.getparent().remove(el)


def _replace_bullets_keeping_format(text_frame, bullets: list[str], body: str = "") -> None:
    if not bullets and not body:
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = ""
        return
    items = bullets if bullets else [body]
    paragraphs = list(text_frame.paragraphs)
    if not paragraphs:
        text_frame.text = items[0]
        for extra in items[1:]:
            text_frame.add_paragraph().text = extra
        return
    first_para = paragraphs[0]
    runs = list(first_para.runs)
    if runs:
        runs[0].text = items[0]
        for run in runs[1:]:
            run.text = ""
    else:
        first_para.add_run().text = items[0]
    for para in paragraphs[1:]:
        para._p.getparent().remove(para._p)
    for extra in items[1:]:
        new_p = deepcopy(first_para._p)
        first_para._p.addnext(new_p)
        for child in new_p.iter():
            if child.tag.split("}")[-1] == "r":
                t_elements = [item for item in child.iter() if item.tag.split("}")[-1] == "t"]
                if t_elements:
                    t_elements[0].text = extra
                    for t in t_elements[1:]:
                        t.text = ""
                break
    container = first_para._p.getparent()
    siblings_after_first = []
    found_first = False
    for child in list(container):
        if found_first and child.tag.split("}")[-1] == "p":
            siblings_after_first.append(child)
        if child is first_para._p:
            found_first = True
    if len(siblings_after_first) == len(items) - 1:
        for sib in siblings_after_first:
            container.remove(sib)
        for sib in reversed(siblings_after_first):
            first_para._p.addnext(sib)


def _fill_body_slots(body_shapes: list, bullets: list[str], subtitle: str, body_text: str) -> None:
    if not body_shapes:
        return
    slots = [shape for shape in body_shapes if shape is not None and shape.has_text_frame]
    if not slots:
        return
    fallback_text = subtitle or body_text or ""
    if len(slots) == 1:
        only = slots[0]
        if bullets:
            _replace_bullets_keeping_format(only.text_frame, bullets, body=body_text)
        elif fallback_text:
            _replace_text_keeping_format(only.text_frame, fallback_text)
        else:
            _replace_text_keeping_format(only.text_frame, "")
        return
    if not bullets:
        _replace_text_keeping_format(slots[0].text_frame, fallback_text)
        for shape in slots[1:]:
            _replace_text_keeping_format(shape.text_frame, "")
        return
    n_slots = len(slots)
    n_bullets = len(bullets)
    if n_bullets >= n_slots:
        for idx, shape in enumerate(slots):
            if idx < n_slots - 1:
                _replace_text_keeping_format(shape.text_frame, bullets[idx])
            else:
                tail = bullets[n_slots - 1 :]
                if len(tail) == 1:
                    _replace_text_keeping_format(shape.text_frame, tail[0])
                else:
                    _replace_bullets_keeping_format(shape.text_frame, tail)
    else:
        for idx, shape in enumerate(slots):
            if idx < n_bullets:
                _replace_text_keeping_format(shape.text_frame, bullets[idx])
            else:
                _replace_text_keeping_format(shape.text_frame, "")


def _delete_slide(presentation: Presentation, index: int) -> None:
    xml_slides = presentation.slides._sldIdLst  # noqa: SLF001
    slides_list = list(xml_slides)
    target = slides_list[index]
    relationship_id = target.attrib["{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"]
    presentation.part.drop_rel(relationship_id)
    xml_slides.remove(target)


def _duplicate_slide(presentation: Presentation, source_index: int) -> int:
    source = presentation.slides[source_index]
    blank = presentation.slide_layouts[6] if len(presentation.slide_layouts) > 6 else presentation.slide_layouts[0]
    dest = presentation.slides.add_slide(blank)
    for shape in list(dest.shapes):
        node = shape._element  # noqa: SLF001
        node.getparent().remove(node)
    for shape in source.shapes:
        new_el = deepcopy(shape._element)  # noqa: SLF001
        dest.shapes._spTree.insert_element_before(new_el, "p:extLst")  # noqa: SLF001
    return len(presentation.slides) - 1


def _slot_render_slide(slide, slot_data: dict, layout_slides: list, template_idx: int | None) -> None:
    if template_idx is None or template_idx >= len(layout_slides):
        _legacy_render_slide(slide, {"bullets": list(slot_data.values())}, layout_slides, None)
        return
    layout_page = layout_slides[template_idx]
    elements = layout_page.get("elements") or []
    shapes = list(slide.shapes)
    has_typed = any(element.get("text_type") for element in elements)
    if not has_typed:
        values = [value for value in slot_data.values() if value]
        _legacy_render_slide(slide, {"bullets": values}, layout_slides, template_idx)
        return
    for el_idx, element in enumerate(elements):
        if el_idx >= len(shapes):
            break
        shape = shapes[el_idx]
        if not shape.has_text_frame:
            continue
        element_id = element.get("id")
        text_type = element.get("text_type") or "noop"
        if element_id in slot_data and slot_data[element_id]:
            new_text = slot_data[element_id]
            _replace_text_keeping_format(shape.text_frame, new_text)
            _shrink_font_if_overflow(shape, new_text)
            _enable_text_autofit(shape.text_frame)
            continue
        current = shape.text_frame.text or ""
        if _is_placeholder_text(current):
            _replace_text_keeping_format(shape.text_frame, "")
            continue
        if text_type in {"title", "subtitle", "content", "item", "itemTitle"}:
            _replace_text_keeping_format(shape.text_frame, "")
    for shape in slide.shapes:
        if shape.has_text_frame and _is_placeholder_text(shape.text_frame.text or ""):
            _replace_text_keeping_format(shape.text_frame, "")


def _legacy_render_slide(slide, slide_data: dict, layout_slides: list, template_idx: int | None) -> None:
    role_map: dict[int, str] | None = None
    if template_idx is not None and template_idx < len(layout_slides):
        elements = layout_slides[template_idx].get("elements") or []
        local_map: dict[int, str] = {}
        for idx, element in enumerate(elements):
            role = element.get("role") or element.get("text_type")
            if role:
                local_map[idx] = role
        role_map = local_map or None

    title_shape, body_shapes, leftover_shapes = _classify_slide_shapes(slide, role_by_index=role_map)
    title_text = (slide_data.get("title") or "").strip()
    subtitle_text = (slide_data.get("subtitle") or "").strip()
    bullets = [str(item).strip() for item in (slide_data.get("bullets") or []) if str(item).strip()]
    body_text = (slide_data.get("body") or "").strip()
    if title_shape is not None and title_shape.has_text_frame and title_text:
        _replace_text_keeping_format(title_shape.text_frame, title_text)
    if body_shapes:
        _fill_body_slots(body_shapes, bullets, subtitle_text, body_text)
    for shape in leftover_shapes:
        if shape.has_text_frame and _is_placeholder_text(shape.text_frame.text or ""):
            _replace_text_keeping_format(shape.text_frame, "")


def render_pptx_from_slides(
    pptx_template_bytes: bytes,
    slides_payload: list[dict],
    *,
    overall_title: str = "",
    template_layout: dict | None = None,
) -> bytes:
    if not slides_payload:
        raise ValueError("slides 不能为空。")
    presentation = Presentation(io.BytesIO(pptx_template_bytes))
    initial_count = len(presentation.slides)
    if initial_count == 0:
        raise ValueError("模板内不包含任何幻灯片。")
    layout_slides = (template_layout or {}).get("slides") or []
    target_indices: list[int] = []
    for idx, slide_payload in enumerate(slides_payload):
        target_idx = slide_payload.get("template_page_idx", idx)
        try:
            target_idx = int(target_idx)
        except (TypeError, ValueError):
            target_idx = idx
        target_indices.append(max(0, min(target_idx, initial_count - 1)))
    for target_idx in target_indices:
        _duplicate_slide(presentation, target_idx)
    for _ in range(initial_count):
        _delete_slide(presentation, 0)
    for idx, slide_payload in enumerate(slides_payload):
        slide = presentation.slides[idx]
        template_idx = target_indices[idx]
        slot_data = slide_payload.get("slot_data") or {}
        if not slot_data and ("title" in slide_payload or "bullets" in slide_payload):
            _legacy_render_slide(slide, slide_payload, layout_slides, template_idx)
        else:
            _slot_render_slide(slide, slot_data, layout_slides, template_idx)
        note_text = (slide_payload.get("speaker_note") or "").strip()
        if note_text:
            try:
                slide.notes_slide.notes_text_frame.text = note_text
            except Exception:
                pass
    if overall_title:
        try:
            presentation.core_properties.title = overall_title
        except Exception:
            pass
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()
