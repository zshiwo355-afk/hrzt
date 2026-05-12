"""PPTX 最小生成：结构化 JSON -> python-pptx 演示文稿。"""
from __future__ import annotations

import json
import re
import site
import sys
import uuid
from io import BytesIO
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
except ModuleNotFoundError:
    user_site = site.getusersitepackages()
    if user_site and user_site not in sys.path:
        sys.path.append(user_site)
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

from app.config import OSS_UPLOAD_PREFIX, UPLOAD_META_DIR, UPLOADS_DIR, oss_configured
from app.providers import oss as oss_provider
from app.models import AttachmentLite, ChatRequest, HistoryMessageItem
from app.services import model_capability_service
from app.services import chat_service
from app.services.attachment_service import (
    build_attachment_public,
    materialized_attachment_file,
    materialize_attachment_file,
    read_attachment_meta_any,
    save_attachment_meta_db,
)
from app.services import ppt_template_renderer
from app.services import rag_service, ofox_gemini_search_service, ofox_responses_service
from app.storage import safe_json_dump


def _history_items_from_task(task_data: dict) -> list[HistoryMessageItem]:
    items: list[HistoryMessageItem] = []
    for row in task_data.get("history_messages") or []:
        if not isinstance(row, dict):
            continue
        attachments = []
        for att in row.get("attachments") or []:
            if not isinstance(att, dict):
                continue
            attachments.append(
                AttachmentLite(
                    id=str(att.get("id") or ""),
                    name=str(att.get("name") or ""),
                    category=str(att.get("category") or "document"),
                )
            )
        items.append(
            HistoryMessageItem(
                role=str(row.get("role") or "user"),
                text=str(row.get("text") or ""),
                attachments=attachments,
                has_image_result=bool(row.get("has_image_result")),
            )
        )
    return items


def _extract_json_object(text: str) -> dict:
    raw = str(text or "").strip()
    if raw.startswith("```"):
        raw = re.sub(r"^```(?:json)?\s*", "", raw, flags=re.IGNORECASE)
        raw = re.sub(r"\s*```$", "", raw)
    start = raw.find("{")
    end = raw.rfind("}")
    if start >= 0 and end > start:
        raw = raw[start : end + 1]
    parsed = json.loads(raw)
    if not isinstance(parsed, dict):
        raise ValueError("模型没有返回可解析的 PPT JSON 对象")
    return parsed


def _build_ppt_messages(task_data: dict, *, knowledge_context: str = "") -> list[dict]:
    system_prompt = (
        "你是 PPT 结构生成器。"
        "请根据用户需求输出严格 JSON，不要解释，不要 Markdown 代码块。"
        '返回 schema: {"title": str, "subtitle": str, "theme_hint": str, "slides": ['
        '{"type": "cover|agenda|content|table|image|ending", "title": str, "subtitle": str, '
        '"bullets": [str], "paragraphs": [str], '
        '"table": {"columns": [str], "rows": [[Any]]}, '
        '"image": {"title": str, "caption": str, "image_hint": str}, '
        '"notes": str}'
        "]}. "
        "要求："
        "1. 至少返回 3 页，最多 8 页；"
        "2. 第一页优先 cover，最后一页优先 ending；"
        "3. 如果适合列表展示，放到 bullets；"
        "4. 如果适合表格展示，使用 table slide；"
        "5. 若用户上传了图片且内容适合配图，允许输出 image slide；"
        "6. 不要输出样式、坐标、颜色值。"
    )
    current_user_content = chat_service.build_current_user_message_content(
        str(task_data.get("prompt") or ""),
        task_data.get("attachment_ids") or [],
    )
    return chat_service.build_text_chat_messages(
        system_prompt=system_prompt,
        summary=str(task_data.get("summary") or ""),
        history_messages=_history_items_from_task(task_data),
        current_user_content=current_user_content,
        knowledge_context=knowledge_context,
    )


def _build_artifact_rag_bundle(task_data: dict) -> dict:
    if not bool(task_data.get("use_rag")):
        return {
            "used": False,
            "query": "",
            "sources": [],
            "context_text": "",
            "note": "",
            "rag_status": "",
        }
    req = ChatRequest(
        conversation_id=int(task_data.get("conversation_id") or 0),
        model=str(task_data.get("model") or ""),
        reasoning_mode=str(task_data.get("reasoning_mode") or "") or None,
        prompt=str(task_data.get("prompt") or ""),
        use_rag=bool(task_data.get("use_rag")),
        rag_query="",
        use_web_search=False,
        attachment_ids=list(task_data.get("attachment_ids") or []),
        summary=str(task_data.get("summary") or ""),
        history_messages=_history_items_from_task(task_data),
    )
    try:
        return rag_service.build_rag_bundle(req)
    except Exception:
        return {
            "used": False,
            "query": str(task_data.get("prompt") or "").strip(),
            "sources": [],
            "context_text": "",
            "note": "知识库服务暂不可用",
            "rag_status": "error",
        }


def _merge_artifact_sources(rag_sources: list[dict], web_sources: list[dict]) -> list[dict]:
    merged: list[dict] = []
    for row in rag_sources or []:
        if isinstance(row, dict):
            item = dict(row)
            item["source_type"] = "rag"
            merged.append(item)
    for row in web_sources or []:
        if isinstance(row, dict):
            item = dict(row)
            item["source_type"] = "web"
            merged.append(item)
    return merged


def _generate_ppt_raw(task_data: dict) -> tuple[str, str, list[dict], str]:
    rag_bundle = _build_artifact_rag_bundle(task_data)
    messages = _build_ppt_messages(task_data, knowledge_context=rag_bundle.get("context_text") or "")
    model = str(task_data.get("model") or "")
    adapter = model_capability_service.build_text_request_adapter(
        model,
        str(task_data.get("reasoning_mode") or "") or None,
    )
    note = chat_service.merge_response_notes(rag_bundle.get("note") or "")
    sources = _merge_artifact_sources(rag_bundle.get("sources") or [], [])
    rag_status = rag_bundle.get("rag_status") or ""

    if bool(task_data.get("use_web_search")) and chat_service.supports_builtin_web_search(model):
        try:
            parsed = ofox_responses_service.parse_responses_web_search_payload(
                ofox_responses_service.call_responses_api(
                    model=model,
                    input_payload=messages,
                    tools=ofox_responses_service.build_responses_web_search_tool(),
                    reasoning_effort=str(
                        (adapter.get("chat_completions_extra") or {}).get("reasoning_effort") or ""
                    ) or None,
                )
            )
            return (
                str(parsed.get("content") or ""),
                chat_service.merge_response_notes(note, "已使用模型内建联网搜索"),
                _merge_artifact_sources(rag_bundle.get("sources") or [], parsed.get("sources") or []),
                rag_status,
            )
        except Exception:
            note = chat_service.merge_response_notes(note, "模型内建联网搜索暂不可用，已按普通文件生成回答")
    elif bool(task_data.get("use_web_search")) and chat_service.supports_gemini_builtin_web_search(model):
        try:
            parsed = ofox_gemini_search_service.parse_gemini_search_payload(
                ofox_gemini_search_service.call_gemini_search_api(
                    model=model,
                    input_payload=messages,
                    reasoning_mode=str(task_data.get("reasoning_mode") or "") or None,
                )
            )
            return (
                str(parsed.get("content") or ""),
                chat_service.merge_response_notes(note, "已使用模型内建联网搜索"),
                _merge_artifact_sources(rag_bundle.get("sources") or [], parsed.get("sources") or []),
                rag_status,
            )
        except Exception:
            note = chat_service.merge_response_notes(note, "模型内建联网搜索暂不可用，已按普通文件生成回答")
    elif bool(task_data.get("use_web_search")):
        note = chat_service.merge_response_notes(note, chat_service.web_search_fallback_note(model))

    raw = chat_service.call_chat_completion(
        model=model,
        messages=messages,
        reasoning_mode=str(task_data.get("reasoning_mode") or "") or None,
        temperature=0.2,
    )
    return raw, note, sources, rag_status


def _normalize_ppt_schema(data: dict) -> dict:
    title = str(data.get("title") or "演示文稿").strip() or "演示文稿"
    subtitle = str(data.get("subtitle") or "").strip()
    theme_hint = str(data.get("theme_hint") or "").strip()
    slides: list[dict[str, Any]] = []
    for idx, row in enumerate(data.get("slides") or []):
        if not isinstance(row, dict):
            continue
        slide_type = str(row.get("type") or "content").strip().lower()
        if slide_type not in {"cover", "agenda", "content", "table", "image", "ending"}:
            slide_type = "content"
        slides.append(
            {
                "type": slide_type,
                "title": str(row.get("title") or f"第{idx + 1}页").strip() or f"第{idx + 1}页",
                "subtitle": str(row.get("subtitle") or "").strip(),
                "bullets": [str(item).strip() for item in (row.get("bullets") or []) if str(item).strip()],
                "paragraphs": [str(item).strip() for item in (row.get("paragraphs") or []) if str(item).strip()],
                "table": row.get("table") if isinstance(row.get("table"), dict) else None,
                "image": row.get("image") if isinstance(row.get("image"), dict) else None,
                "notes": str(row.get("notes") or "").strip(),
            }
        )
    if not slides:
        slides = [
            {"type": "cover", "title": title, "subtitle": subtitle, "bullets": [], "paragraphs": [], "table": None, "image": None, "notes": ""},
            {"type": "content", "title": "内容概览", "subtitle": "", "bullets": ["请补充内容"], "paragraphs": [], "table": None, "image": None, "notes": ""},
            {"type": "ending", "title": "感谢聆听", "subtitle": "", "bullets": [], "paragraphs": [], "table": None, "image": None, "notes": ""},
        ]
    if slides[0]["type"] != "cover":
        slides.insert(0, {"type": "cover", "title": title, "subtitle": subtitle, "bullets": [], "paragraphs": [], "table": None, "image": None, "notes": ""})
    if slides[-1]["type"] != "ending":
        slides.append({"type": "ending", "title": "感谢聆听", "subtitle": "", "bullets": [], "paragraphs": [], "table": None, "image": None, "notes": ""})
    return {"title": title, "subtitle": subtitle, "theme_hint": theme_hint, "slides": slides[:8]}


def _safe_file_stem(value: str, fallback: str) -> str:
    text = re.sub(r'[\\/:*?"<>|]+', "_", str(value or "").strip())
    text = re.sub(r"\s+", " ", text).strip().strip(".")
    if not text:
        text = fallback
    return text[:60]


def _attachment_image_paths(task_data: dict) -> list[Path]:
    paths: list[Path] = []
    for attachment_id in task_data.get("attachment_ids") or []:
        meta = read_attachment_meta_any(str(attachment_id))
        if not meta or meta.get("category") != "image":
            continue
        path = materialize_attachment_file(meta)
        if path and path.is_file():
            paths.append(path)
    return paths


def _apply_run_style(run, *, size: int, bold: bool = False, color: tuple[int, int, int] = (31, 31, 31)) -> None:
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


def _add_title(slide, text: str, *, top: float = 0.5, size: int = 24) -> None:
    shape = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.2), Inches(0.7))
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    _apply_run_style(run, size=size, bold=True, color=(21, 54, 92))


def _add_subtitle(slide, text: str, *, top: float = 1.3, size: int = 12) -> None:
    if not text:
        return
    shape = slide.shapes.add_textbox(Inches(0.75), Inches(top), Inches(10.8), Inches(0.45))
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _apply_run_style(run, size=size, color=(88, 102, 119))


def _add_bullets(slide, bullets: list[str], *, top: float = 1.8, height: float = 4.5) -> None:
    shape = slide.shapes.add_textbox(Inches(0.9), Inches(top), Inches(10.0), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.clear()
    for idx, bullet in enumerate(bullets or []):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = bullet
        _apply_run_style(run, size=18)


def _add_paragraphs(slide, paragraphs: list[str], *, top: float = 1.8, height: float = 4.5) -> None:
    shape = slide.shapes.add_textbox(Inches(0.9), Inches(top), Inches(10.0), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.clear()
    for idx, line in enumerate(paragraphs or []):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        _apply_run_style(run, size=18)


def _add_table(slide, title: str, table_data: dict) -> None:
    columns = [str(col).strip() for col in (table_data.get("columns") or []) if str(col).strip()]
    rows = [row for row in (table_data.get("rows") or []) if isinstance(row, list)]
    _add_title(slide, title)
    if not columns:
        _add_paragraphs(slide, ["未提供表格列定义。"])
        return
    row_count = max(1, len(rows) + 1)
    shape = slide.shapes.add_table(row_count, len(columns), Inches(0.7), Inches(1.5), Inches(11.0), Inches(4.8))
    table = shape.table
    for col_idx, name in enumerate(columns):
        cell = table.cell(0, col_idx)
        cell.text = name
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                _apply_run_style(run, size=12, bold=True, color=(255, 255, 255))
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(31, 78, 121)
    max_rows = min(len(rows), 8)
    for r_idx in range(max_rows):
        values = list(rows[r_idx]) + [""] * max(0, len(columns) - len(rows[r_idx]))
        for c_idx, value in enumerate(values[: len(columns)]):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(value)
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(247, 251, 255)


def _add_image_slide(slide, title: str, caption: str, image_path: Path | None) -> None:
    _add_title(slide, title)
    if image_path and image_path.is_file():
        slide.shapes.add_picture(str(image_path), Inches(1.1), Inches(1.5), width=Inches(8.8), height=Inches(4.7))
    else:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(1.7), Inches(8.8), Inches(3.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(236, 243, 250)
        shape.line.color.rgb = RGBColor(180, 196, 214)
        tf = shape.text_frame
        tf.text = "未提供可用图片，当前页使用占位展示"
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                _apply_run_style(run, size=18, color=(88, 102, 119))
    _add_subtitle(slide, caption, top=6.45, size=12)


def _render_pptx_file(schema: dict, task_data: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    image_metas = []
    for attachment_id in task_data.get("attachment_ids") or []:
        meta = read_attachment_meta_any(str(attachment_id))
        if meta and meta.get("category") == "image":
            image_metas.append(meta)
    image_iter = iter(image_metas)

    for slide_data in schema["slides"]:
        slide = prs.slides.add_slide(blank_layout)
        slide_type = slide_data["type"]
        if slide_type == "cover":
            _add_title(slide, slide_data["title"] or schema["title"], top=1.3, size=28)
            _add_subtitle(slide, slide_data.get("subtitle") or schema.get("subtitle") or "", top=2.4, size=18)
        elif slide_type == "agenda":
            bullets = slide_data.get("bullets") or slide_data.get("paragraphs") or []
            _add_title(slide, slide_data["title"] or "目录")
            _add_bullets(slide, bullets, top=1.8, height=4.6)
        elif slide_type == "table":
            _add_table(slide, slide_data["title"], slide_data.get("table") or {})
        elif slide_type == "image":
            image_info = slide_data.get("image") or {}
            caption = str(image_info.get("caption") or slide_data.get("subtitle") or "").strip()
            image_meta = next(image_iter, None)
            with materialized_attachment_file(image_meta) as image_path:
                _add_image_slide(slide, slide_data["title"], caption, image_path)
        elif slide_type == "ending":
            _add_title(slide, slide_data["title"] or "感谢聆听", top=2.4, size=28)
            _add_subtitle(slide, slide_data.get("subtitle") or "欢迎交流", top=3.4, size=18)
        else:
            _add_title(slide, slide_data["title"])
            bullets = slide_data.get("bullets") or []
            paragraphs = slide_data.get("paragraphs") or []
            if bullets:
                _add_bullets(slide, bullets, top=1.8, height=4.6)
            elif paragraphs:
                _add_paragraphs(slide, paragraphs, top=1.8, height=4.6)
            else:
                _add_paragraphs(slide, [slide_data.get("notes") or "请补充内容"], top=1.8, height=4.6)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def _render_pptx_with_best_strategy(schema: dict, task_data: dict) -> tuple[bytes, str]:
    if ppt_template_renderer.is_template_eligible(schema):
        try:
            content, meta = ppt_template_renderer.render_with_builtin_template(
                schema,
                prompt=str(task_data.get("prompt") or ""),
            )
            return content, f"template:{meta.get('template_label') or meta.get('template_id') or 'builtin'}"
        except Exception:
            pass
    return _render_pptx_file(schema, task_data), "fixed_layout"


def run_pptx_task(task_data: dict) -> dict:
    raw, generation_note, generation_sources, rag_status = _generate_ppt_raw(task_data)
    parsed = _extract_json_object(raw)
    schema = _normalize_ppt_schema(parsed)
    filename = _safe_file_stem(schema["title"], "演示文稿") + ".pptx"
    content, render_mode = _render_pptx_with_best_strategy(schema, task_data)

    attachment_id = uuid.uuid4().hex
    stored_name = f"{attachment_id}.pptx"
    stored_path = UPLOADS_DIR / stored_name
    storage = "local"
    oss_key = ""
    stored_path.write_bytes(content)
    if oss_configured():
        oss_key = f"{OSS_UPLOAD_PREFIX}/generated/{attachment_id}.pptx"
        if oss_provider.write_bytes(
            oss_key,
            content,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ):
            storage = "oss"
    meta = {
        "id": attachment_id,
        "original_name": filename,
        "stored_name": stored_name,
        "stored_path": str(stored_path),
        "storage": storage,
        "oss_key": oss_key if storage == "oss" else "",
        "suffix": ".pptx",
        "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "category": "document",
        "size": len(content or b""),
    }
    safe_json_dump(UPLOAD_META_DIR / f"{attachment_id}.json", meta)
    save_attachment_meta_db(meta, task_id=str(task_data.get("id") or "") or None, source="generated_artifact")
    attachment = build_attachment_public(meta)
    link = f"/api/attachments/{attachment['id']}"
    note = chat_service.merge_response_notes(generation_note, "任务已完成，已按演示文稿模板渲染。")
    if render_mode == "fixed_layout":
        note = chat_service.merge_response_notes(generation_note, "任务已完成，当前内容包含表格或图片页，已按兼容布局渲染。")
    return {
        "text": f"已生成PPT演示文稿：[下载 {attachment['name']}]({link})",
        "note": note,
        "artifact_type": "pptx",
        "attachment": attachment,
        "sources": generation_sources,
        "rag_status": rag_status,
    }
